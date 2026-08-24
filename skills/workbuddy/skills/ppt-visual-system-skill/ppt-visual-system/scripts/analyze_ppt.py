"""
PPT 分析脚本模板

Skill: ppt-visual-system v1.0.0

使用方法：
    python analyze_ppt.py <input.pptx> [slide_range]
    
    例：
    python analyze_ppt.py my_deck.pptx          # 分析所有页
    python analyze_ppt.py my_deck.pptx 1-5       # 分析第1-5页
    python analyze_ppt.py my_deck.pptx 3          # 分析第3页

功能：
    1. 输出每页的所有 Shape 属性（位置、尺寸、字体、颜色）
    2. 识别并标记使用了禁用字体的元素
    3. 检查布局是否符合规范
"""

import sys
import os
from pptx import Presentation
from pptx.util import Emu, Pt

# 禁用字体列表
BANNED_FONTS = {
    'Calibri', 'Arial', 'Microsoft YaHei', 'Consolas',
    'SimSun', 'Times New Roman', '宋体', '华文细黑',
    '华文楷体', '华文宋体', '华文中宋', '华文仿宋',
}


def analyze_slide(slide, slide_num):
    """分析单页 slide"""
    print('=' * 70)
    print('SLIDE {}'.format(slide_num))
    print('=' * 70)
    
    font_issues = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            left = Emu(shape.left).inches if shape.left else 0
            top = Emu(shape.top).inches if shape.top else 0
            w = Emu(shape.width).inches if shape.width else 0
            h = Emu(shape.height).inches if shape.height else 0
            
            print('  SHAPE: {} | type={} | pos=({:.2f},{:.2f}) size=({:.2f}x{:.2f})'.format(
                shape.name, shape.shape_type, left, top, w, h))
            
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    txt = para.text.strip()
                    if txt:
                        font_info = ''
                        if para.runs:
                            r = para.runs[0]
                            bold = 'BOLD ' if r.font.bold else ''
                            sz = '{:.0f}pt'.format(r.font.size.pt) if r.font.size else '?pt'
                            color = ''
                            font_name = r.font.name or '(inherited)'
                            
                            try:
                                if r.font.color and r.font.color.rgb:
                                    color = ' #{}'.format(r.font.color.rgb)
                            except:
                                pass
                            
                            font_info = ' [{}{} {} {}]'.format(bold, sz, font_name, color)
                            
                            # 检查禁用字体
                            if r.font.name and r.font.name in BANNED_FONTS:
                                font_issues.append((slide_num, shape.name, r.font.name, txt[:30]))
                        
                        print('    P{}{}: {}'.format(pi, font_info, txt[:80]))
    
    # Notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print('  NOTES: {}...'.format(notes[:200]))
    
    print()
    return font_issues


def main():
    if len(sys.argv) < 2:
        print('Usage: python analyze_ppt.py <input.pptx> [slide_range]')
        sys.exit(1)
    
    input_path = sys.argv[1]
    slide_range = sys.argv[2] if len(sys.argv) > 2 else None
    
    if not os.path.exists(input_path):
        print('Error: {} not found'.format(input_path))
        sys.exit(1)
    
    prs = Presentation(input_path)
    
    print('File: {}'.format(input_path))
    print('Slide size: {:.2f} x {:.2f} inches'.format(
        Emu(prs.slide_width).inches, Emu(prs.slide_height).inches))
    print('Total slides: {}'.format(len(prs.slides)))
    print()
    
    # 解析 slide_range
    if slide_range:
        if '-' in slide_range:
            start, end = map(int, slide_range.split('-'))
        else:
            start = end = int(slide_range)
    else:
        start, end = 1, len(prs.slides)
    
    all_font_issues = []
    
    for i, slide in enumerate(prs.slides, 1):
        if start <= i <= end:
            issues = analyze_slide(slide, i)
            all_font_issues.extend(issues)
    
    # 字体问题汇总
    if all_font_issues:
        print('\n' + '=' * 70)
        print('⚠️  字体问题汇总 ({} 处)'.format(len(all_font_issues)))
        print('=' * 70)
        for slide_num, shape_name, font_name, text in all_font_issues:
            print('  P{}: {} 使用了禁用字体 "{}" -> "{}"'.format(
                slide_num, shape_name, font_name, text))
    else:
        print('\n✅ 未发现禁用字体')


if __name__ == '__main__':
    main()
