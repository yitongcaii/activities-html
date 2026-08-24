"""
AISee Security Logo - Generate .pptx file
渐变流光风格 + 盾牌安全元素
"""
import math
import random
import io
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = r"d:\AI\workbuddy\20260403\logo_assets"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ============================================================
# Color definitions
# ============================================================
DEEP_BG = (10, 14, 28)
CORE_BLUE = (60, 140, 230)
CYAN = (100, 210, 255)
PURPLE = (140, 100, 255)
LIGHT_BLUE = (150, 200, 255)
WHITE = (240, 245, 255)

# ============================================================
# Drawing helpers
# ============================================================

def draw_shield_points(cx, cy, size):
    """Return shield polygon points."""
    top = cy - size
    bottom = cy + int(size * 1.3)
    mid_w = int(size * 0.85)
    top_w = int(size * 0.65)
    return [
        (cx, top),
        (cx + top_w, cy - int(size*0.3)),
        (cx + mid_w, cy + int(size*0.2)),
        (cx + int(mid_w*0.7), cy + int(size*0.7)),
        (cx, bottom),
        (cx - int(mid_w*0.7), cy + int(size*0.7)),
        (cx - mid_w, cy + int(size*0.2)),
        (cx - top_w, cy - int(size*0.3)),
    ]

def draw_shield(draw, cx, cy, size, fill, outline=None, width=2):
    pts = draw_shield_points(cx, cy, size)
    draw.polygon(pts, fill=fill, outline=outline, width=width)

def draw_glow_ellipse(draw, cx, cy, rx, ry, color, layers=8):
    for i in range(layers, 0, -1):
        alpha = int(40 / i)
        c = color[:3] + (alpha,) if len(color) == 3 else color
        draw.ellipse([cx-rx-i*12, cy-ry-i*12, cx+rx+i*12, cy+ry+i*12], fill=c)

def draw_circuit(draw, cx, cy, size, color):
    random.seed(42)
    for _ in range(10):
        a = random.uniform(0, 2*math.pi)
        r1 = size * random.uniform(0.15, 0.45)
        r2 = size * random.uniform(0.4, 0.85)
        x1 = cx + int(r1*math.cos(a))
        y1 = cy + int(r1*math.sin(a))
        x2 = cx + int(r2*math.cos(a + random.uniform(-0.4, 0.4)))
        y2 = cy + int(r2*math.sin(a + random.uniform(-0.4, 0.4)))
        draw.line([(x1,y1),(x2,y2)], fill=color, width=1)
        draw.ellipse([x2-2,y2-2,x2+2,y2+2], fill=color)

def draw_eye(draw, cx, cy, size, color):
    ew = int(size * 0.55)
    eh = int(size * 0.22)
    for off in range(4, 0, -1):
        alpha = 60 + (4-off)*45
        ec = color[:3] + (alpha,) if len(color) <= 3 else color
        draw.arc([cx-ew-off, cy-eh-off, cx+ew+off, cy+eh+off], 200, 340, fill=ec, width=2)
        draw.arc([cx-ew-off, cy-eh-off, cx+ew+off, cy+eh+off], 20, 160, fill=ec, width=2)
    draw.ellipse([cx-9,cy-9,cx+9,cy+9], fill=color)
    draw.ellipse([cx-4,cy-4,cx+4,cy+4], fill=(255,255,255,230))

def draw_particles(img, cx, cy, radius, count=35, color=(80,180,255)):
    draw = ImageDraw.Draw(img)
    random.seed(123)
    for _ in range(count):
        a = random.uniform(0, 2*math.pi)
        d = random.uniform(radius*0.7, radius*1.6)
        x = cx + int(d*math.cos(a))
        y = cy + int(d*math.sin(a))
        s = random.uniform(1, 3.5)
        alpha = random.randint(60, 200)
        draw.ellipse([x-s,y-s,x+s,y+s], fill=color+(alpha,))

def draw_streaks(img, cx, cy, radius, count=15):
    draw = ImageDraw.Draw(img)
    random.seed(77)
    for _ in range(count):
        a = random.uniform(0, 2*math.pi)
        length = random.uniform(radius*0.4, radius*1.3)
        sr = radius * 0.65
        x1 = cx + int(sr*math.cos(a))
        y1 = cy + int(sr*math.sin(a))
        x2 = cx + int((sr+length)*math.cos(a))
        y2 = cy + int((sr+length)*math.sin(a))
        for i in range(25):
            t = i/25.0
            px = int(x1+(x2-x1)*t)
            py = int(y1+(y2-y1)*t)
            alpha = int(160*(1-t))
            c = (80+int(60*t), 170+int(80*t), 255, alpha)
            draw.ellipse([px-1,py-1,px+1,py+1], fill=c)

def draw_glow_line(draw, x1, y1, x2, y2, color, width=2):
    for i in range(6, 0, -1):
        alpha = int(40/i*2)
        c = color[:3]+(alpha,)
        draw.line([(x1,y1),(x2,y2)], fill=c, width=width+i*3)
    draw.line([(x1,y1),(x2,y2)], fill=color, width=width)

# ============================================================
# Logo generators
# ============================================================

def create_logo_dark(size=800):
    """深色背景 + 流光盾牌"""
    w, h = size, size
    img = Image.new('RGBA', (w,h), DEEP_BG + (255,))
    
    # BG gradient
    bg = Image.new('RGBA', (w,h), (0,0,0,0))
    bgd = ImageDraw.Draw(bg)
    for y in range(h):
        t = y/h
        r = int(DEEP_BG[0]*(1-t) + 20*t)
        g = int(DEEP_BG[1]*(1-t) + 30*t)
        b = int(DEEP_BG[2]*(1-t) + 60*t)
        bgd.line([(0,y),(w,y)], fill=(r,g,b,255))
    img = Image.alpha_composite(img, bg)

    draw = ImageDraw.Draw(img)
    cx, cy = w//2, h//2 - 25
    ss = 140

    # Outer glow
    draw_glow_ellipse(draw, cx, cy, ss+20, ss+20, (50,120,255), 10)

    # Shield
    sl = Image.new('RGBA', (w,h), (0,0,0,0))
    sd = ImageDraw.Draw(sl)
    draw_shield(sd, cx, cy, ss+3, (30,70,170,180), (90,170,255,210), 3)
    draw_shield(sd, cx, cy, ss-8, (50,110,200,70), None, 0)
    draw_circuit(sd, cx, cy, ss, (110,190,255,140))
    draw_eye(sd, cx, cy-12, ss, (100,220,255,230))
    img = Image.alpha_composite(img, sl)

    draw_particles(img, cx, cy, ss, 45, (70,170,255))
    draw_streaks(img, cx, cy, ss, 16)

    # Glow line
    gl = Image.new('RGBA', (w,h), (0,0,0,0))
    gd = ImageDraw.Draw(gl)
    ly = cy + int(ss*1.5)
    draw_glow_line(gd, cx-190, ly, cx+190, ly, (100,190,255,200))
    img = Image.alpha_composite(img, gl)

    # Text
    try:
        fl = ImageFont.truetype("C:/Windows/Fonts/arialbd.ttf", 70)
        fs = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 26)
    except:
        fl = ImageFont.load_default()
        fs = ImageFont.load_default()

    tl = Image.new('RGBA', (w,h), (0,0,0,0))
    td = ImageDraw.Draw(tl)
    ty = ly + 22
    for i in range(6,0,-1):
        td.text((cx-148, ty-2), "AISee", font=fl, fill=(60,140,255,int(25/i*2)))
    td.text((cx-148, ty), "AI", font=fl, fill=(90,200,255,255))
    td.text((cx-22, ty), "See", font=fl, fill=(160,130,255,255))
    td.text((cx-105, ty+78), "AI-Powered Security", font=fs, fill=(140,170,210,170))
    img = Image.alpha_composite(img, tl)

    return img.convert('RGB')


def create_logo_light(size=800):
    """白色背景 + 蓝紫渐变盾牌"""
    w, h = size, size
    img = Image.new('RGBA', (w,h), (255,255,255,255))
    draw = ImageDraw.Draw(img)
    cx, cy = w//2, h//2 - 35
    ss = 130

    # Shadow
    sh = Image.new('RGBA', (w,h), (0,0,0,0))
    shd = ImageDraw.Draw(sh)
    draw_shield(shd, cx+5, cy+5, ss, (0,0,0,25))
    img = Image.alpha_composite(img, sh)

    # Shield
    sl = Image.new('RGBA', (w,h), (0,0,0,0))
    sd = ImageDraw.Draw(sl)
    draw_shield(sd, cx, cy, ss+2, (35,65,165,210), (70,140,230,240), 3)
    
    # Gradient overlay
    go = Image.new('RGBA', (w,h), (0,0,0,0))
    god = ImageDraw.Draw(go)
    pts = draw_shield_points(cx, cy, ss-6)
    for y in range(min(p[1] for p in pts), max(p[1] for p in pts)):
        t = (y - min(p[1] for p in pts)) / max(1, max(p[1] for p in pts) - min(p[1] for p in pts))
        r = int(40 + t*70)
        g = int(80 + t*30)
        b = int(200 + t*40)
        a = int(90*(1-abs(t-0.5)*2))
        god.line([(cx-ss,y),(cx+ss,y)], fill=(r,g,b,a))
    
    mask = Image.new('L', (w,h), 0)
    md = ImageDraw.Draw(mask)
    md.polygon(pts, fill=255)
    go = Image.composite(go, sl, mask)
    sl = Image.alpha_composite(sl, go)
    
    draw_circuit(sd, cx, cy, ss, (90,150,240,120))
    draw_eye(sd, cx, cy-8, ss, (70,150,240,235))
    img = Image.alpha_composite(img, sl)

    draw_particles(img, cx, cy, ss, 28, (70,140,240))
    draw_streaks(img, cx, cy, ss, 12)

    gl = Image.new('RGBA', (w,h), (0,0,0,0))
    gd = ImageDraw.Draw(gl)
    ly = cy + int(ss*1.5)
    draw_glow_line(gd, cx-175, ly, cx+175, ly, (70,140,235,160))
    img = Image.alpha_composite(img, gl)

    try:
        fl = ImageFont.truetype("C:/Windows/Fonts/arialbd.ttf", 66)
        fs = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 24)
    except:
        fl = ImageFont.load_default()
        fs = ImageFont.load_default()

    tl = Image.new('RGBA', (w,h), (0,0,0,0))
    td = ImageDraw.Draw(tl)
    ty = ly + 18
    td.text((cx-138, ty), "AI", font=fl, fill=(35,95,200,255))
    td.text((cx-12, ty), "See", font=fl, fill=(115,75,195,255))
    td.text((cx-100, ty+72), "Intelligent Security", font=fs, fill=(100,120,165,180))
    img = Image.alpha_composite(img, tl)

    return img.convert('RGB')


def create_logo_aurora(size=800):
    """极光背景 + 增强流光盾牌"""
    w, h = size, size
    img = Image.new('RGBA', (w,h), (8,10,25,255))
    
    # BG gradient
    bg = Image.new('RGBA', (w,h), (0,0,0,0))
    bgd = ImageDraw.Draw(bg)
    for y in range(h):
        t = y/h
        r = int(8 + 18*math.sin(t*math.pi))
        g = int(10 + 25*math.sin(t*math.pi*1.3))
        b = int(25 + 35*math.sin(t*math.pi*0.7))
        bgd.line([(0,y),(w,y)], fill=(r,g,b,255))
    img = Image.alpha_composite(img, bg)

    # Aurora blobs
    au = Image.new('RGBA', (w,h), (0,0,0,0))
    aud = ImageDraw.Draw(au)
    random.seed(999)
    for _ in range(60):
        x = random.randint(0,w)
        y = random.randint(0,h//2)
        s = random.uniform(20,90)
        alpha = random.randint(3,12)
        c = random.choice([(40,80,200,alpha),(80,40,180,alpha),(60,120,200,alpha)])
        aud.ellipse([x-s,y-s,x+s,y+s], fill=c)
    img = Image.alpha_composite(img, au)

    draw = ImageDraw.Draw(img)
    cx, cy = w//2, h//2 - 15
    ss = 145

    draw_glow_ellipse(draw, cx, cy, ss+25, ss+25, (55,100,220), 12)

    sl = Image.new('RGBA', (w,h), (0,0,0,0))
    sd = ImageDraw.Draw(sl)
    draw_shield(sd, cx, cy, ss+3, (22,55,125,160), (95,175,255,190), 3)
    draw_shield(sd, cx, cy, ss-9, (40,80,165,75), None, 0)
    
    # Top inner glow
    for i in range(6,0,-1):
        alpha = int(18/i)
        iy = cy - ss + int(ss*0.4)
        sd.ellipse([cx-int(ss*0.3)-i*5, iy-i*5, cx+int(ss*0.3)+i*5, iy+int(ss*0.4)+i*5], fill=(115,195,255,alpha))

    draw_circuit(sd, cx, cy, ss, (105,195,255,135))
    draw_eye(sd, cx, cy-8, ss, (115,215,255,240))
    img = Image.alpha_composite(img, sl)

    draw_particles(img, cx, cy, ss, 55, (75,175,255))
    draw_streaks(img, cx, cy, ss, 20)

    # Gradient glow line (cyan to purple)
    gl = Image.new('RGBA', (w,h), (0,0,0,0))
    gd = ImageDraw.Draw(gl)
    ly = cy + int(ss*1.5)
    steps = 120
    for i in range(steps):
        x = cx - 200 + int(400*i/steps)
        t = i/steps
        r = int(55+t*85)
        g = int(140-t*45)
        b = int(255-t*30)
        for j in range(5,0,-1):
            alpha = int(35/j)
            gd.line([(x,ly-j*2),(x,ly+j*2)], fill=(r,g,b,alpha), width=1)
    gd.line([(cx-200,ly),(cx+200,ly)], fill=(145,175,255,220), width=2)
    img = Image.alpha_composite(img, gl)

    try:
        fl = ImageFont.truetype("C:/Windows/Fonts/arialbd.ttf", 72)
        fs = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 28)
    except:
        fl = ImageFont.load_default()
        fs = ImageFont.load_default()

    tl = Image.new('RGBA', (w,h), (0,0,0,0))
    td = ImageDraw.Draw(tl)
    ty = ly + 22
    for i in range(7,0,-1):
        td.text((cx-148, ty-3), "AISee", font=fl, fill=(75,145,255,int(20/i*2)))
    td.text((cx-148, ty), "AI", font=fl, fill=(75,195,255,255))
    td.text((cx-22, ty), "See", font=fl, fill=(155,115,255,255))
    td.text((cx-115, ty+80), "AI  \u00b7  Security  \u00b7  Intelligence", font=fs, fill=(135,165,210,160))
    img = Image.alpha_composite(img, tl)

    return img.convert('RGB')


def create_logo_icon(size=600):
    """图标版本（仅盾牌+眼）"""
    w, h = size, size
    img = Image.new('RGBA', (w,h), (10,14,30,255))
    
    bg = Image.new('RGBA', (w,h), (0,0,0,0))
    bgd = ImageDraw.Draw(bg)
    for y in range(h):
        t = y/h
        r = int(10+8*math.sin(t*math.pi))
        g = int(14+12*math.sin(t*math.pi))
        b = int(30+20*math.sin(t*math.pi))
        bgd.line([(0,y),(w,y)], fill=(r,g,b,255))
    img = Image.alpha_composite(img, bg)

    draw = ImageDraw.Draw(img)
    cx, cy = w//2, h//2
    ss = 150

    draw_glow_ellipse(draw, cx, cy, ss+22, ss+22, (50,110,230), 10)

    sl = Image.new('RGBA', (w,h), (0,0,0,0))
    sd = ImageDraw.Draw(sl)
    draw_shield(sd, cx, cy, ss+2, (25,60,140,165), (85,170,255,205), 3)
    draw_shield(sd, cx, cy, ss-10, (42,88,180,75), None, 0)
    draw_circuit(sd, cx, cy, ss, (105,185,255,140))
    draw_eye(sd, cx, cy-5, ss, (105,215,255,235))
    img = Image.alpha_composite(img, sl)

    draw_particles(img, cx, cy, ss, 50, (65,165,255))
    draw_streaks(img, cx, cy, ss, 18)

    return img.convert('RGB')


def create_bg_gradient(w=1920, h=1080, color1=(8,10,25), color2=(18,28,55)):
    """Create a dark gradient background for slides."""
    img = Image.new('RGB', (w, h))
    draw = ImageDraw.Draw(img)
    for y in range(h):
        t = y / h
        r = int(color1[0]*(1-t) + color2[0]*t)
        g = int(color1[1]*(1-t) + color2[1]*t)
        b = int(color1[2]*(1-t) + color2[2]*t)
        draw.line([(0,y),(w,y)], fill=(r,g,b))
    return img


# ============================================================
# Save images
# ============================================================
print("Generating logo images...")

logo_dark = create_logo_dark(800)
logo_light = create_logo_light(800)
logo_aurora = create_logo_aurora(800)
logo_icon = create_logo_icon(600)

logo_dark_path = os.path.join(OUTPUT_DIR, "logo_dark.png")
logo_light_path = os.path.join(OUTPUT_DIR, "logo_light.png")
logo_aurora_path = os.path.join(OUTPUT_DIR, "logo_aurora.png")
logo_icon_path = os.path.join(OUTPUT_DIR, "logo_icon.png")

logo_dark.save(logo_dark_path)
logo_light.save(logo_light_path)
logo_aurora.save(logo_aurora_path)
logo_icon.save(logo_icon_path)

# Smaller versions for comparison
logo_dark_sm = logo_dark.resize((400,400), Image.LANCZOS)
logo_light_sm = logo_light.resize((400,400), Image.LANCZOS)
logo_aurora_sm = logo_aurora.resize((400,400), Image.LANCZOS)
logo_dark_sm_path = os.path.join(OUTPUT_DIR, "logo_dark_sm.png")
logo_light_sm_path = os.path.join(OUTPUT_DIR, "logo_light_sm.png")
logo_aurora_sm_path = os.path.join(OUTPUT_DIR, "logo_aurora_sm.png")
logo_dark_sm.save(logo_dark_sm_path)
logo_light_sm.save(logo_light_sm_path)
logo_aurora_sm.save(logo_aurora_sm_path)

# Favicon
favicon = logo_icon.resize((120,120), Image.LANCZOS)
favicon_path = os.path.join(OUTPUT_DIR, "favicon.png")
favicon.save(favicon_path)

# Slide backgrounds
bg_main = create_bg_gradient(1920, 1080, (8,10,25), (18,28,55))
bg_cover = create_bg_gradient(1920, 1080, (6,8,20), (15,20,45))
bg_end = create_bg_gradient(1920, 1080, (8,10,25), (22,18,48))
bg_main_path = os.path.join(OUTPUT_DIR, "bg_main.png")
bg_cover_path = os.path.join(OUTPUT_DIR, "bg_cover.png")
bg_end_path = os.path.join(OUTPUT_DIR, "bg_end.png")
bg_main.save(bg_main_path)
bg_light_bg = create_bg_gradient(1920, 1080, (240,242,248), (220,225,238))
bg_light_path = os.path.join(OUTPUT_DIR, "bg_light.png")
bg_light_bg.save(bg_light_path)
bg_main.save(bg_cover_path)
bg_main.save(bg_end_path)

print("All images saved.")

# ============================================================
# Create PPTX
# ============================================================
print("Creating PPTX...")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = 13.333
SLIDE_H = 7.5

# Helper: add background image
def set_bg(slide, img_path):
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = RGBColor(*line_color)
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 0: Cover
# ============================================================
slide0 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide0, bg_cover_path)
slide0.shapes.add_picture(logo_icon_path, Inches(SLIDE_W/2 - 1.8), Inches(0.8), Inches(3.6), Inches(3.6))
add_text(slide0, 1.5, 4.6, 10.3, 1.0, "AISee Security", 52, True, CYAN, PP_ALIGN.CENTER, "Arial Black")
add_text(slide0, 1.5, 5.6, 10.3, 0.5, "Brand Logo Design", 22, False, (150,175,210), PP_ALIGN.CENTER)
add_text(slide0, 1.5, 6.3, 10.3, 0.4, "\u6e10\u53d8\u6d41\u5149 \u00b7 \u76fe\u724c\u5b88\u62a4 \u00b7 \u667a\u80fd\u5b89\u5168", 16, False, (120,145,185), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 1: Version A - Dark
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide1, bg_main_path)
slide1.shapes.add_picture(logo_dark_path, Inches(1.5), Inches(0.6), Inches(4.5), Inches(4.5))

add_text(slide1, 7.0, 1.0, 5.5, 0.4, "VERSION A", 14, False, (100,130,170), PP_ALIGN.LEFT)
add_text(slide1, 7.0, 1.5, 5.5, 0.7, "\u6df1\u7a7a\u6d41\u5149", 40, True, CYAN, PP_ALIGN.LEFT, "Arial Black")
add_text(slide1, 7.0, 2.4, 5.5, 0.4, "Deep Space Glow", 18, False, (130,155,195), PP_ALIGN.LEFT)

# Description card
card = add_shape(slide1, MSO_SHAPE.RECTANGLE, 7.0, 3.2, 5.3, 3.2, fill_color=(15,22,48), line_color=(50,90,160), line_width=1)

txBox = slide1.shapes.add_textbox(Inches(7.3), Inches(3.5), Inches(4.8), Inches(2.8))
tf = txBox.text_frame
tf.word_wrap = True

items = [
    ("\u2022  \u6df1\u8272\u80cc\u666f\u642d\u914d\u7535\u5149\u84dd\u6d41\u5149\u6548\u679c\uff0c\u79d1\u6280\u611f\u5f3a\u70c8", 15),
    ("", 8),
    ("\u2022  \u76fe\u724c\u5185\u90e8\u5d4c\u5165 AI \u4e4b\u773c\u7b26\u53f7\uff0c\u8c61\u5f81\u667a\u80fd\u6d1e\u5bdf", 15),
    ("", 8),
    ("\u2022  \u7c92\u5b50\u5149\u6548\u73af\u7ed5\u76fe\u724c\uff0c\u8c61\u5f81\u5168\u65b9\u4f4d\u667a\u80fd\u9632\u62a4", 15),
    ("", 8),
    ("\u2022  \u9002\u5408\u5b89\u5168\u4ea7\u54c1\u4e3b\u54c1\u724c\u6807\u8bc6\u3001\u53d1\u5e03\u4f1a\u7b49\u573a\u666f", 15),
]
for i, (text, size) in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(170,190,220)
    p.font.name = "Arial"

# ============================================================
# SLIDE 2: Version B - Light
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, bg_light_path)
slide2.shapes.add_picture(logo_light_path, Inches(1.5), Inches(0.6), Inches(4.5), Inches(4.5))

add_text(slide2, 7.0, 1.0, 5.5, 0.4, "VERSION B", 14, True, (80,70,140), PP_ALIGN.LEFT)
add_text(slide2, 7.0, 1.5, 5.5, 0.7, "\u6781\u7b80\u767d\u5883", 40, True, PURPLE, PP_ALIGN.LEFT, "Arial Black")
add_text(slide2, 7.0, 2.4, 5.5, 0.4, "Minimal White", 18, False, (100,85,155), PP_ALIGN.LEFT)

card = add_shape(slide2, MSO_SHAPE.RECTANGLE, 7.0, 3.2, 5.3, 3.2, fill_color=(248,248,255), line_color=(130,120,190), line_width=1)

txBox = slide2.shapes.add_textbox(Inches(7.3), Inches(3.5), Inches(4.8), Inches(2.8))
tf = txBox.text_frame
tf.word_wrap = True

items2 = [
    ("\u2022  \u767d\u8272\u80cc\u666f\u4e0a\u7684\u84dd\u7d2b\u6e10\u53d8\u76fe\u724c\uff0c\u7b80\u6d01\u4f18\u96c5", 15),
    ("", 8),
    ("\u2022  \u67d4\u548c\u9634\u5f71\u4e0e\u5149\u7c92\u5b50\u4fdd\u6301\u54c1\u724c\u4e00\u81f4\u6027", 15),
    ("", 8),
    ("\u2022  \u9002\u5408\u6d45\u8272\u4e3b\u9898\u6587\u6863\u3001\u5408\u4f5c\u4f19\u4f34\u5c55\u793a", 15),
    ("", 8),
    ("\u2022  \u6253\u5370\u53cb\u597d\uff0c\u9002\u5408\u7eb8\u8d28\u7269\u6599\u4e0e\u5468\u8fb9\u5236\u4f5c", 15),
]
for i, (text, size) in enumerate(items2):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(60,55,90)
    p.font.name = "Arial"

# ============================================================
# SLIDE 3: Version C - Aurora
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3, bg_end_path)
slide3.shapes.add_picture(logo_aurora_path, Inches(1.5), Inches(0.6), Inches(4.5), Inches(4.5))

add_text(slide3, 7.0, 1.0, 5.5, 0.4, "VERSION C", 14, False, (100,130,170), PP_ALIGN.LEFT)
add_text(slide3, 7.0, 1.5, 5.5, 0.7, "\u6781\u5149\u5e7b\u5883", 40, True, (130,180,255), PP_ALIGN.LEFT, "Arial Black")
add_text(slide3, 7.0, 2.4, 5.5, 0.4, "Aurora Fantasy", 18, False, (130,155,195), PP_ALIGN.LEFT)

card = add_shape(slide3, MSO_SHAPE.RECTANGLE, 7.0, 3.2, 5.3, 3.2, fill_color=(15,20,45), line_color=(70,110,190), line_width=1)

txBox = slide3.shapes.add_textbox(Inches(7.3), Inches(3.5), Inches(4.8), Inches(2.8))
tf = txBox.text_frame
tf.word_wrap = True

items3 = [
    ("\u2022  \u6781\u5149\u80cc\u666f\u6548\u679c\u642d\u914d\u589e\u5f3a\u7248\u6d41\u5149\u76fe\u724c", 15),
    ("", 8),
    ("\u2022  \u84dd\u7d2b\u6e10\u53d8\u5149\u5e26\u4ece\u76fe\u724c\u5411\u5916\u8f90\u5c04\uff0c\u89c6\u89c9\u51b2\u51fb\u529b\u6700\u5f3a", 15),
    ("", 8),
    ("\u2022  \u589e\u5f3a\u7684\u7c92\u5b50\u6548\u679c\u4e0e\u5149\u7ebf\u8f90\u5c04", 15),
    ("", 8),
    ("\u2022  \u9002\u5408\u54c1\u724c\u5c55\u793a\u3001\u53d1\u5e03\u4f1a\u3001\u5ba3\u4f20\u6d77\u62a5\u7b49\u9ad8\u5149\u573a\u666f", 15),
]
for i, (text, size) in enumerate(items3):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(170,190,220)
    p.font.name = "Arial"

# ============================================================
# SLIDE 4: Comparison
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide4, bg_main_path)

add_text(slide4, 1, 0.4, 11.3, 0.7, "\u4e09\u7248\u5bf9\u6bd4", 40, True, CYAN, PP_ALIGN.CENTER, "Arial Black")
add_text(slide4, 1, 1.1, 11.3, 0.4, "COMPARISON", 16, False, (120,145,185), PP_ALIGN.CENTER)

# Three cards
card_data = [
    (logo_dark_sm_path, "\u6df1\u7a7a\u6d41\u5149", "\u6df1\u8272\u4e3b\u9898 \u00b7 \u7535\u5149\u84dd \u00b7 \u79d1\u6280\u611f", (15,22,48), (50,90,160)),
    (logo_light_sm_path, "\u6781\u7b80\u767d\u5883", "\u6d45\u8272\u4e3b\u9898 \u00b7 \u84dd\u7d2b\u6e10\u53d8 \u00b7 \u4f18\u96c5", (248,248,255), (130,120,190)),
    (logo_aurora_sm_path, "\u6781\u5149\u5e7b\u5883", "\u6781\u5149\u6548\u679c \u00b7 \u5168\u5f69\u6d41\u5149 \u00b7 \u9707\u64bc", (15,20,45), (70,110,190)),
]

card_w = 3.4
gap = 0.55
start_x = (SLIDE_W - 3*card_w - 2*gap) / 2

for i, (img_path, title, desc, bg_c, border_c) in enumerate(card_data):
    x = start_x + i*(card_w + gap)
    # Card background
    add_shape(slide4, MSO_SHAPE.RECTANGLE, x, 2.0, card_w, 4.8, fill_color=bg_c, line_color=border_c, line_width=1)
    # Logo
    slide4.shapes.add_picture(img_path, x + (card_w-3.0)/2, 2.3, Inches(3.0), Inches(3.0))
    # Title
    add_text(slide4, x, 5.5, card_w, 0.5, title, 22, True, CYAN if i != 1 else PURPLE, PP_ALIGN.CENTER)
    # Description
    add_text(slide4, x, 6.0, card_w, 0.4, desc, 12, False, (130,150,180) if i != 1 else (80,70,120), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Design Elements
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide5, bg_main_path)

add_text(slide5, 1, 0.4, 11.3, 0.7, "\u8bbe\u8ba1\u5143\u7d20", 40, True, WHITE, PP_ALIGN.CENTER, "Arial Black")
add_text(slide5, 1, 1.1, 11.3, 0.4, "DESIGN ELEMENTS", 16, False, (120,145,185), PP_ALIGN.CENTER)

# Icon + description on left, elements grid on right
slide5.shapes.add_picture(logo_icon_path, Inches(1.2), Inches(1.8), Inches(2.8), Inches(2.8))

add_text(slide5, 1.2, 4.8, 2.8, 0.4, "\u56fe\u6807\u7248\u672c", 18, True, (180,200,230), PP_ALIGN.CENTER)
add_text(slide5, 1.2, 5.2, 2.8, 0.3, "App Icon / Favicon", 12, False, (120,145,175), PP_ALIGN.CENTER)

# Four element cards
elements = [
    ("\u76fe\u724c Shield", "\u5b89\u5168\u5b88\u62a4\u7684\u6838\u5fc3\u7b26\u53f7\uff0c\u51e0\u4f55\u5316\u7684\u672a\u6765\u611f\u8bbe\u8ba1", (50,90,160)),
    ("AI \u4e4b\u773c Eye", "\u5d4c\u5165\u76fe\u724c\u5185\u90e8\uff0c\u8c61\u5f81 AI \u667a\u80fd\u6d1e\u5bdf\u529b", (50,90,160)),
    ("\u6d41\u5149\u7c92\u5b50 Glow", "\u73af\u7ed5\u76fe\u724c\u7684\u5149\u7c92\u5b50\u6548\u679c\uff0c\u589e\u5f3a\u52a8\u6001\u79d1\u6280\u611f", (50,90,160)),
    ("\u6e10\u53d8\u8272 Gradient", "\u84dd\u7d2b\u6e10\u53d8\u4e3b\u8272\u8c03\uff0c\u878d\u5408\u51b7\u8272\u79d1\u6280\u4e0e\u6696\u8272\u4fe1\u8d56", (50,90,160)),
]

elem_w = 2.3
elem_gap = 0.35
elem_start_x = 4.8
elem_start_y = 1.8

for i, (title, desc, border_c) in enumerate(elements):
    row = i // 2
    col = i % 2
    x = elem_start_x + col * (elem_w + elem_gap)
    y = elem_start_y + row * 2.5
    
    add_shape(slide5, MSO_SHAPE.RECTANGLE, x, y, elem_w, 2.0, fill_color=(15,22,48), line_color=border_c, line_width=1)
    add_text(slide5, x+0.15, y+0.25, elem_w-0.3, 0.5, title, 18, True, CYAN, PP_ALIGN.LEFT)
    
    txBox = slide5.shapes.add_textbox(Inches(x+0.15), Inches(y+0.85), Inches(elem_w-0.3), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(155,170,200)
    p.font.name = "Arial"

# ============================================================
# SLIDE 6: Application Scenarios
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide6, bg_main_path)

add_text(slide6, 1, 0.4, 11.3, 0.7, "\u5e94\u7528\u573a\u666f", 40, True, WHITE, PP_ALIGN.CENTER, "Arial Black")
add_text(slide6, 1, 1.1, 11.3, 0.4, "APPLICATION SCENARIOS", 16, False, (120,145,185), PP_ALIGN.CENTER)

apps = [
    ("App \u56fe\u6807", "\u79fb\u52a8\u7aef\u4e0e\u684c\u9762\u5e94\u7528"),
    ("Favicon", "\u6d4f\u89c8\u5668\u6807\u7b7e\u9875\u6807\u8bc6"),
    ("\u793e\u4ea4\u5a92\u4f53", "\u5934\u50cf\u4e0e\u54c1\u724c\u4e3b\u9875"),
    ("\u6587\u6863\u6c34\u5370", "\u62a5\u544a\u4e0e\u6587\u6863\u54c1\u724c\u6807\u8bc6"),
    ("\u54c1\u724c\u5468\u8fb9", "T\u6064\u3001\u676f\u5b50\u3001\u5de5\u724c\u7b49"),
    ("\u53d1\u5e03\u4f1a", "Keynote \u4e0e\u5ba3\u4f20\u6d77\u62a5"),
]

app_w = 3.5
app_h = 2.5
app_gap_x = 0.6
app_gap_y = 0.5
app_start_x = (SLIDE_W - 3*app_w - 2*app_gap_x) / 2
app_start_y = 1.9

for i, (title, desc) in enumerate(apps):
    row = i // 3
    col = i % 3
    x = app_start_x + col * (app_w + app_gap_x)
    y = app_start_y + row * (app_h + app_gap_y)
    
    add_shape(slide6, MSO_SHAPE.RECTANGLE, x, y, app_w, app_h, fill_color=(15,22,48), line_color=(45,75,140), line_width=1)
    
    # Small icon
    slide6.shapes.add_picture(favicon_path, x + (app_w-0.8)/2, y+0.2, Inches(0.8), Inches(0.8))
    
    add_text(slide6, x, y+1.2, app_w, 0.5, title, 16, True, (190,210,240), PP_ALIGN.CENTER)
    add_text(slide6, x, y+1.7, app_w, 0.4, desc, 12, False, (130,150,185), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: Thank You
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide7, bg_end_path)
slide7.shapes.add_picture(logo_icon_path, Inches(SLIDE_W/2 - 1.6), Inches(1.2), Inches(3.2), Inches(3.2))
add_text(slide7, 1.5, 4.8, 10.3, 1.0, "Thank You", 52, True, CYAN, PP_ALIGN.CENTER, "Arial Black")
add_text(slide7, 1.5, 5.9, 10.3, 0.5, "AISee \u00b7 \u5b89\u5168\u770b\u5f97\u89c1", 22, False, (140,165,205), PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = r"d:\AI\workbuddy\20260403\AISee_Security_Logo.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
