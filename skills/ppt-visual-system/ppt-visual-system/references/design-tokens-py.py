"""
PPT Design Tokens — python-pptx (Python) 版本

Skill: ppt-visual-system v1.0.0

使用方法：
    from design_tokens_py import *
    
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    solid_fill_bg(sl, DEEP_BLUE)
    tb = add_textbox(sl, Inches(1.2), Inches(2.0), Inches(11), Inches(1.2))
    set_text(tb, '标题', FONT_BOLD, Pt(48), True, WHITE)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════
# Design Tokens — 色彩系统
# ═══════════════════════════════════════════════════════════

# 背景色
DEEP_BLUE    = RGBColor(0x0F, 0x17, 0x2A)  # bgDark — 封面/分隔/结尾
NEAR_WHITE   = RGBColor(0xF8, 0xFA, 0xFC)  # bgLight — 内容页
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)   # 卡片背景/反白文字

# 主题色
TENCENT_BLUE = RGBColor(0x00, 0x52, 0xD9)  # primary — 核心强调
LIGHT_BLUE   = RGBColor(0x05, 0x94, 0xFA)  # primaryLight — 次要强调

# 辅助色
GREEN_ACCENT = RGBColor(0x61, 0xDD, 0xAA)  # green — 正面/成功
RED_ACCENT   = RGBColor(0xDC, 0x26, 0x26)  # red — 警告/重要
YELLOW_ACC   = RGBColor(0xF6, 0xC0, 0x22)  # yellow — 注意/预警
NAVY_BLUE    = RGBColor(0x2B, 0x5F, 0xD9)  # navyBlue — 深蓝辅助

# 文字色
DARK_TEXT    = RGBColor(0x0F, 0x17, 0x2A)  # 最深文字（=DEEP_BLUE）
HEADING_TXT  = RGBColor(0x1E, 0x29, 0x3B)  # textDark — 主标题
BODY_TEXT    = RGBColor(0x64, 0x74, 0x8B)  # textGray — 正文/描述
SLATE_400    = RGBColor(0x94, 0xA3, 0xB8)  # textLight — 辅助说明
MUTED_NUM    = RGBColor(0xE2, 0xE8, 0xF0)  # textMuted — 装饰数字

# 边框色
LIGHT_BORDER = RGBColor(0xE2, 0xE8, 0xF0)  # border

# 特殊用途
HIGHLIGHT_BG = RGBColor(0xEF, 0xF6, 0xFF)  # 高亮条背景
WARNING_BG   = RGBColor(0xFE, 0xF2, 0xF2)  # 警告标签背景

# ═══════════════════════════════════════════════════════════
# 字体系统
# ═══════════════════════════════════════════════════════════
FONT_BOLD = 'TencentSans W7'   # 标题、金句、标签
FONT_REG  = 'TencentSans W3'   # 正文、描述、注释

# ═══════════════════════════════════════════════════════════
# 尺寸常量
# ═══════════════════════════════════════════════════════════
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ═══════════════════════════════════════════════════════════
# 辅助函数
# ═══════════════════════════════════════════════════════════

def set_font(run, name=None, size=None, bold=False, color=None):
    """设置 run 的字体属性"""
    f = run.font
    if name:
        f.name = name
    if size:
        f.size = size
    f.bold = bold
    if color:
        f.color.rgb = color


def add_textbox(slide, left, top, width, height):
    """添加文本框并返回 shape"""
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(shape, text, font_name=FONT_REG, font_size=Pt(12),
             bold=False, color=BODY_TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    """设置 shape 的单段文字"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    for p in tf.paragraphs:
        p.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=font_size, bold=bold, color=color)
    return shape


def add_para(text_frame, text, font_name=FONT_REG, font_size=Pt(12),
             bold=False, color=BODY_TEXT, align=PP_ALIGN.LEFT,
             space_before=None, space_after=None):
    """向已有 text_frame 追加一个段落"""
    p = text_frame.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=font_size, bold=bold, color=color)
    return p


def solid_fill_bg(slide, color):
    """设置 slide 纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════
# 页面工厂函数
# ═══════════════════════════════════════════════════════════

def create_presentation():
    """创建标准演示文稿"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def cover_slide(prs, title, subtitle=None, speaker=None):
    """创建封面页"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    solid_fill_bg(sl, DEEP_BLUE)
    
    # 顶部装饰条
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.08), TENCENT_BLUE)
    
    # 标题
    tb = add_textbox(sl, Inches(1.2), Inches(2.0), Inches(11), Inches(1.2))
    set_text(tb, title, FONT_BOLD, Pt(48), True, WHITE, PP_ALIGN.LEFT)
    
    # 副标题
    if subtitle:
        tb = add_textbox(sl, Inches(1.2), Inches(3.3), Inches(10), Inches(0.8))
        set_text(tb, subtitle, FONT_REG, Pt(22), False, SLATE_400)
    
    # 署名
    if speaker:
        add_rect(sl, Inches(1.2), Inches(5.5), Inches(2), Inches(0.04), TENCENT_BLUE)
        tb = add_textbox(sl, Inches(1.2), Inches(5.8), Inches(6), Inches(0.6))
        set_text(tb, speaker, FONT_REG, Pt(11), False, SLATE_400)
    
    return sl


def divider_slide(prs, part_num, title, subtitle=None):
    """创建章节分隔页"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    solid_fill_bg(sl, NEAR_WHITE)
    
    # 大淡色编号
    tb = add_textbox(sl, Inches(1.2), Inches(1.5), Inches(3), Inches(1.5))
    set_text(tb, f'{part_num:02d}', FONT_BOLD, Pt(72), True, MUTED_NUM)
    
    # 标题
    tb = add_textbox(sl, Inches(1.2), Inches(3.0), Inches(10), Inches(1.0))
    set_text(tb, title, FONT_BOLD, Pt(42), True, DARK_TEXT)
    
    # 副标题
    if subtitle:
        tb = add_textbox(sl, Inches(1.2), Inches(4.2), Inches(8), Inches(0.6))
        set_text(tb, subtitle, FONT_REG, Pt(18), False, BODY_TEXT)
    
    return sl


def content_slide(prs, section_tag, title):
    """创建标准内容页（带 section tag 和标题）"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    solid_fill_bg(sl, NEAR_WHITE)
    
    # Section tag
    tb = add_textbox(sl, Inches(0.8), Inches(0.4), Inches(3), Inches(0.3))
    set_text(tb, section_tag, FONT_BOLD, Pt(10), True, TENCENT_BLUE)
    
    # Title
    tb = add_textbox(sl, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7))
    set_text(tb, title, FONT_BOLD, Pt(30), True, DARK_TEXT)
    
    return sl


def ending_slide(prs, main_text, sub_text=None, speaker=None):
    """创建结尾页"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    solid_fill_bg(sl, DEEP_BLUE)
    
    # 顶部装饰条
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.08), TENCENT_BLUE)
    
    # 主文案
    tb = add_textbox(sl, Inches(1.2), Inches(2.0), Inches(11), Inches(1.2))
    set_text(tb, main_text, FONT_BOLD, Pt(44), True, WHITE)
    
    # 副文案
    if sub_text:
        tb = add_textbox(sl, Inches(1.2), Inches(3.5), Inches(10), Inches(1.5))
        set_text(tb, sub_text, FONT_REG, Pt(22), False, SLATE_400)
    
    # Q&A
    add_rect(sl, Inches(5.5), Inches(5.3), Inches(2.5), Inches(0.04), TENCENT_BLUE)
    tb = add_textbox(sl, Inches(5.5), Inches(5.5), Inches(2.5), Inches(0.5))
    set_text(tb, 'Q & A', FONT_REG, Pt(16), False, SLATE_400, PP_ALIGN.CENTER)
    
    # 署名
    if speaker:
        tb = add_textbox(sl, Inches(4.5), Inches(6.2), Inches(4.5), Inches(0.6))
        set_text(tb, speaker, FONT_REG, Pt(11), False, SLATE_400, PP_ALIGN.CENTER)
    
    return sl
